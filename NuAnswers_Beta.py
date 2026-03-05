import streamlit as st
from openai import OpenAI
import pandas as pd
from datetime import datetime, timezone, timedelta
import os
import tempfile
from pathlib import Path
import PyPDF2
import docx
import pptx
import csv
import xlrd
import openpyxl
import io
import base64
from zoneinfo import ZoneInfo
import re

# Set page config
st.set_page_config(
    page_title="NuAnswers",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Hide all default Streamlit elements we don't want to show
st.markdown("""
    <style>
        /* Hide page names in sidebar */
        span.css-10trblm.e16nr0p30 {
            display: none;
        }
        /* Hide the default Streamlit menu button */
        button.css-1rs6os.edgvbvh3 {
            display: none;
        }
        /* Hide "streamlit app" text */
        .css-17ziqus {
            display: none;
        }
        /* Hide development mode indicator */
        .stDeployButton {
            display: none;
        }
        /* Sidebar styling */
        .css-1d391kg {
            padding-top: 1rem;
        }
    </style>
""", unsafe_allow_html=True)

# Tutoring hours configuration (in 24-hour format)
TUTORING_HOURS = {
    "Monday": [("9:00", "11:00")],    # 9:00 AM - 11:00 AM
    "Tuesday": [("9:00", "11:00")],   # 9:00 AM - 11:00 AM
    "Wednesday": [("9:00", "11:00")], # 9:00 AM - 11:00 AM
    "Thursday": [("9:00", "11:00")],  # 9:00 AM - 11:00 AM
    "Friday": [("9:00", "11:00")],    # 9:00 AM - 11:00 AM
}

def is_within_tutoring_hours():
    """Check if current time is within tutoring hours."""
    # Get current time in Eastern Time
    et_tz = ZoneInfo("America/New_York")
    current_time = datetime.now(et_tz)
    
    current_day = current_time.strftime("%A")
    current_hour = current_time.hour
    current_minute = current_time.minute
    current_time_float = current_hour + (current_minute / 60)
    
    # Debug logging with more detailed time information
    st.session_state.debug_time = {
        "current_time": current_time.strftime("%I:%M %p"),
        "current_time_24h": current_time.strftime("%H:%M"),
        "current_day": current_day,
        "current_hour": current_hour,
        "current_minute": current_minute,
        "current_time_float": current_time_float,
        "timezone": "Eastern Time (ET)"
    }
    
    # If it's not a tutoring day, return False
    if current_day not in TUTORING_HOURS:
        st.session_state.debug_time["reason"] = "Not a tutoring day"
        return False
    
    # Check each tutoring time slot
    for start_time_str, end_time_str in TUTORING_HOURS[current_day]:
        # Convert time strings to float hours (e.g., "13:30" -> 13.5)
        start_hour, start_minute = map(int, start_time_str.split(":"))
        end_hour, end_minute = map(int, end_time_str.split(":"))
        
        start_time_float = start_hour + (start_minute / 60)
        end_time_float = end_hour + (end_minute / 60)
        
        # Add to debug info
        st.session_state.debug_time.update({
            "tutoring_start": f"{start_hour:02d}:{start_minute:02d}",
            "tutoring_end": f"{end_hour:02d}:{end_minute:02d}",
            "start_time_float": start_time_float,
            "end_time_float": end_time_float,
            "is_within": current_time_float >= start_time_float and current_time_float <= end_time_float
        })
        
        # Check if current time is within the tutoring slot
        if current_time_float >= start_time_float and current_time_float <= end_time_float:
            st.session_state.debug_time["reason"] = "Within tutoring hours"
            return True
            
    st.session_state.debug_time["reason"] = "Outside tutoring hours"
    return False

# Configure data directory
DATA_DIR = Path("/data" if os.path.exists("/data") else ".")
REGISTRATION_DATA_PATH = DATA_DIR / "registration_data.csv"
FEEDBACK_DATA_PATH = DATA_DIR / "feedback_data.csv"
TOPIC_DATA_PATH = DATA_DIR / "topic_data.csv"
COMPLETION_DATA_PATH = DATA_DIR / "completion_data.csv"
RESPONSE_TIMES_PATH = DATA_DIR / "response_times.csv"
CONTENT_ACCESS_PATH = DATA_DIR / "content_access.csv"
RESOLUTION_TIMES_PATH = DATA_DIR / "resolution_times.csv"

# Create data directory if it doesn't exist
DATA_DIR.mkdir(exist_ok=True)

# Initialize all session state variables
if "registered" not in st.session_state:
    st.session_state.registered = False
if "start_time" not in st.session_state:
    st.session_state.start_time = None
if "user_data" not in st.session_state:
    st.session_state.user_data = {}
if "messages" not in st.session_state:
    st.session_state.messages = [
        {"role": "assistant", "content": "Hello! I'm NuAnswers. I'm here to help you understand concepts and work through problems. What would you like to work on today?"}
    ]
if "registration_data" not in st.session_state:
    st.session_state.registration_data = pd.DataFrame(columns=[
        "timestamp", "full_name", "student_id", "student_email", "grade", "campus",
        "major", "course_name", "course_id", "professor", "professor_email", "usage_time_minutes"
    ])
if "uploaded_documents" not in st.session_state:
    st.session_state.uploaded_documents = []
if "search_query" not in st.session_state:
    st.session_state.search_query = ""
if "doc_to_delete" not in st.session_state:
    st.session_state.doc_to_delete = None
if "feedback_data" not in st.session_state:
    st.session_state.feedback_data = []
if "topic_data" not in st.session_state:
    st.session_state.topic_data = []
if "completion_data" not in st.session_state:
    st.session_state.completion_data = []
if "logout_initiated" not in st.session_state:
    st.session_state.logout_initiated = False
if "feedback_submitted" not in st.session_state:
    st.session_state.feedback_submitted = False
if "chat_started" not in st.session_state:
    st.session_state.chat_started = False  # True after user submits "Start new chat" (course info)
if "show_returning_lookup" not in st.session_state:
    st.session_state.show_returning_lookup = True  # Show returning user option first
if "response_times" not in st.session_state:
    st.session_state.response_times = []
if "content_access" not in st.session_state:
    st.session_state.content_access = []
if "resolution_times" not in st.session_state:
    st.session_state.resolution_times = []

def save_to_csv(data, filepath):
    """Save data to CSV file with error handling"""
    try:
        # Convert to DataFrame if it's a dictionary
        if isinstance(data, dict):
            df = pd.DataFrame([data])
        elif isinstance(data, list):
            df = pd.DataFrame(data)
        else:
            df = data
            
        # Load existing data if file exists
        if filepath.exists():
            existing_df = pd.read_csv(filepath)
            df = pd.concat([existing_df, df], ignore_index=True)
        
        # Save to CSV
        df.to_csv(filepath, index=False)
    except Exception as e:
        st.error(f"Failed to save data to {filepath}: {str(e)}")

def save_registration(user_data, start_time):
    """Save registration/session data to CSV (full session with usage time)."""
    et_tz = ZoneInfo("America/New_York")
    end_time = datetime.now(et_tz)
    if start_time is None:
        start_time = end_time
    if start_time.tzinfo is None:
        start_time = start_time.replace(tzinfo=et_tz)
    else:
        start_time = start_time.astimezone(et_tz)
    usage_time = (end_time - start_time).total_seconds() / 60
    new_registration = {
        "timestamp": end_time.strftime("%Y-%m-%d %H:%M:%S"),
        "full_name": user_data.get("full_name", ""),
        "student_id": user_data.get("student_id", ""),
        "student_email": user_data.get("student_email", ""),
        "grade": user_data.get("grade", ""),
        "campus": user_data.get("campus", ""),
        "major": user_data.get("major", ""),
        "course_name": user_data.get("course_name", ""),
        "course_id": user_data.get("course_id", ""),
        "professor": user_data.get("professor", ""),
        "professor_email": user_data.get("professor_email", ""),
        "usage_time_minutes": usage_time
    }
    save_to_csv(new_registration, REGISTRATION_DATA_PATH)


def lookup_student_from_csv(student_id, student_email):
    """Return account dict for returning user from CSV, or None if not found."""
    if not REGISTRATION_DATA_PATH.exists():
        return None
    try:
        df = pd.read_csv(REGISTRATION_DATA_PATH)
        if df.empty or "student_id" not in df.columns:
            return None
        email_col = "student_email" if "student_email" in df.columns else "email"
        if email_col not in df.columns:
            return None
        mask = (df["student_id"].astype(str).str.strip() == str(student_id).strip()) & (
            df[email_col].astype(str).str.strip().str.lower() == str(student_email).strip().lower()
        )
        matches = df.loc[mask]
        if matches.empty:
            return None
        row = matches.iloc[-1]
        return {
            "full_name": row.get("full_name", ""),
            "student_id": str(row.get("student_id", "")),
            "student_email": str(row.get(email_col, "")),
            "grade": row.get("grade", ""),
            "campus": row.get("campus", ""),
            "major": row.get("major", ""),
            "course_name": "",
            "course_id": "",
            "professor": "",
            "professor_email": "",
        }
    except Exception:
        return None

def track_content_access(content_id, content_type):
    """Track content access patterns"""
    entry = {
        "timestamp": datetime.now(ZoneInfo("America/New_York")).strftime("%Y-%m-%d %H:%M:%S"),
        "content_id": content_id,
        "content_type": content_type,
        "user_id": st.session_state.user_data.get("full_name")
    }
    st.session_state.content_access.append(entry)
    save_to_csv(entry, CONTENT_ACCESS_PATH)

def track_resolution_time(start_time, end_time, topic):
    """Track problem resolution time"""
    resolution_time = (end_time - start_time).total_seconds() / 60  # Convert to minutes
    entry = {
        "timestamp": end_time.strftime("%Y-%m-%d %H:%M:%S"),
        "resolution_time": resolution_time,
        "topic": topic,
        "user_id": st.session_state.user_data.get("full_name")
    }
    st.session_state.resolution_times.append(entry)
    save_to_csv(entry, RESOLUTION_TIMES_PATH)

# Create a sidebar
with st.sidebar:
    st.title("📚 NuAnswers")

# Initialize session state for form
if "form_major" not in st.session_state:
    st.session_state.form_major = "Accounting"

def update_major():
    """Update the major in session state"""
    st.session_state.form_major = st.session_state.temp_major
    st.rerun()

# ----- Step 1: Not registered — show "Returning user" or "Create account" (6 fields only) -----
if not st.session_state.registered:
    st.title("📝 NuAnswers")
    st.write("DISCLAIMER: All information will be kept private and not shared with anyone. Used only by BAP: Nu Sigma Chapter.")

    # Returning user: quick re-entry with Student ID + Email
    with st.expander("Returning user? Enter your details to continue", expanded=st.session_state.show_returning_lookup):
        with st.form("returning_user_form"):
            return_id = st.text_input("FDU Student ID (7 digits)", key="return_id")
            return_email = st.text_input("FDU Student Email", key="return_email")
            lookup_clicked = st.form_submit_button("Continue")
        if lookup_clicked and return_id and return_email:
            is_valid_id = return_id.isdigit() and len(return_id) == 7
            is_valid_email = return_email.endswith("@student.fdu.edu") or return_email.endswith("@fdu.edu")
            if not is_valid_id:
                st.error("Student ID must be exactly 7 digits.")
            elif not is_valid_email:
                st.error("Email must be a valid FDU email (@student.fdu.edu or @fdu.edu).")
            else:
                account = lookup_student_from_csv(return_id, return_email)
                try:
                    from supabase_db import get_student_by_credentials
                    if account is None:
                        account = get_student_by_credentials(return_id, return_email)
                except Exception:
                    pass
                if account:
                    st.session_state.user_data = account
                    st.session_state.registered = True
                    st.session_state.chat_started = False
                    st.success("Welcome back! Please fill in your course details below to start a new chat.")
                    st.rerun()
                else:
                    st.info("No account found with that Student ID and email. Please create an account below.")

    st.subheader("Create account (first time only)")
    st.caption("Fill out: Full name, FDU Student ID, FDU email, Grade, Campus, Major.")
    with st.form("registration_form", clear_on_submit=False):
        full_name = st.text_input("Full Name")
        student_id = st.text_input("FDU Student ID (7 digits)")
        student_email = st.text_input("FDU Student Email (@student.fdu.edu or @fdu.edu)")
        grade = st.selectbox("Grade", ["Freshman", "Sophomore", "Junior", "Senior", "Graduate"])
        campus = st.selectbox("Campus", ["Florham", "Metro", "Vancouver"])
        major = st.selectbox(
            "Major",
            ["Accounting", "Finance", "MIS [Management Information Systems]"]
        )
        submitted = st.form_submit_button("Create account")
        if submitted:
            is_valid_student_id = student_id.isdigit() and len(student_id) == 7
            is_valid_student_email = student_email.endswith("@student.fdu.edu") or student_email.endswith("@fdu.edu")
            if not all([full_name, student_id, student_email]):
                st.error("Please fill in Full name, Student ID, and Email.")
            elif not is_valid_student_id:
                st.error("Student ID must be exactly 7 digits.")
            elif not is_valid_student_email:
                st.error("Email must be a valid FDU email address.")
            else:
                st.session_state.user_data = {
                    "full_name": full_name,
                    "student_id": student_id,
                    "student_email": student_email,
                    "grade": grade,
                    "campus": campus,
                    "major": major,
                    "course_name": "",
                    "course_id": "",
                    "professor": "",
                    "professor_email": ""
                }
                et_tz = ZoneInfo("America/New_York")
                st.session_state.start_time = datetime.now(et_tz)
                # Save account-only row so returning user lookup works
                save_registration(st.session_state.user_data, st.session_state.start_time)
                try:
                    from supabase_db import save_registration_data
                    save_registration_data(st.session_state.user_data, st.session_state.start_time)
                except Exception:
                    pass
                st.session_state.registered = True
                st.session_state.chat_started = False
                st.success("Account created! Now enter your course details below to start a chat.")
                st.rerun()

# ----- Step 2: Registered but no chat started — "Start new chat" (4 fields only) -----
elif st.session_state.registered and not st.session_state.chat_started:
    st.title("📝 Start a new chat")
    st.write("Enter the course details for this session. You only need to fill this out each time you start a new chat.")
    with st.form("new_chat_form", clear_on_submit=False):
        course_name = st.text_input("Which class are you taking that relates to what you need help in?")
        course_id = st.text_input("Course ID (Format: DEPT_####_##)", help="Examples: ACCT_2021_01, FIN_3250_02")
        professor = st.text_input("Professor's Name")
        professor_email = st.text_input("Professor's Email")
        if course_id:
            valid_prefixes = ['ACCT', 'ECON', 'FIN', 'MIS', 'WMA']
            is_valid_course_id = bool(re.match(f"^({'|'.join(valid_prefixes)})_\\d{{4}}_\\d{{2}}$", course_id))
            if not is_valid_course_id:
                st.error("Use format: ACCT_####_##, ECON_####_##, FIN_####_##, MIS_####_##, or WMA_####_##")
        submitted = st.form_submit_button("Start chat")
        if submitted:
            if not all([course_name, course_id, professor, professor_email]):
                st.error("Please fill in all four fields.")
            elif not re.match(r"^(ACCT|ECON|FIN|MIS|WMA)_\d{4}_\d{2}$", course_id):
                st.error("Please enter a valid Course ID format.")
            else:
                ud = st.session_state.user_data
                st.session_state.user_data = {
                    **ud,
                    "course_name": course_name,
                    "course_id": course_id,
                    "professor": professor,
                    "professor_email": professor_email
                }
                et_tz = ZoneInfo("America/New_York")
                st.session_state.start_time = datetime.now(et_tz)
                st.session_state.chat_started = True
                st.session_state.messages = [
                    {"role": "assistant", "content": "Hello! I'm NuAnswers. I'm here to help you understand concepts and work through problems. What would you like to work on today?"}
                ]
                st.rerun()

# Function to extract text from different file types
def extract_text_from_file(file):
    file_extension = Path(file.name).suffix.lower()
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
        tmp_file.write(file.getvalue())
        tmp_file_path = tmp_file.name
    
    try:
        if file_extension == '.pdf':
            text = extract_text_from_pdf(tmp_file_path)
        elif file_extension == '.docx':
            text = extract_text_from_docx(tmp_file_path)
        elif file_extension == '.txt':
            with open(tmp_file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        elif file_extension == '.pptx':
            text = extract_text_from_pptx(tmp_file_path)
        elif file_extension == '.csv':
            text = extract_text_from_csv(tmp_file_path)
        elif file_extension in ['.xls', '.xlsx']:
            text = extract_text_from_excel(tmp_file_path)
        else:
            st.error(f"Unsupported file type: {file_extension}")
            return None
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return None
    finally:
        os.unlink(tmp_file_path)
    
    return text

def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_csv(file_path):
    text = ""
    with open(file_path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        for row in reader:
            text += ", ".join(row) + "\n"
    return text

def extract_text_from_excel(file_path):
    text = ""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            text += f"\nSheet: {sheet_name}\n"
            text += df.to_string(index=False) + "\n"
    except Exception as e:
        st.error(f"Error reading Excel file: {str(e)}")
        return None
    return text

# Function to search within documents
def search_in_documents(query, documents):
    if not query:
        return documents
    query = query.lower()
    results = []
    for doc in documents:
        if query in doc['name'].lower() or query in doc['content'].lower():
            results.append(doc)
    return results

# Main application logic for registered users who have started a chat (entered course details)
if st.session_state.registered and st.session_state.chat_started:
    # Show the introduction message once at the top
    st.title("💬 NuAnswers")
    st.write(
        "Hello! I am NuAnswers, Beta Alpha Psi: Nu Sigma Chapter's AI Tutor Bot. I'm here to help you understand concepts and work through problems. "
        "Remember, I won't give you direct answers, but I'll guide you to find them yourself. "
        "I can help you with accounting equations, financial ratios, financial statements, and time value of money concepts."
    )
    # Handle logout process with feedback first
    if st.session_state.logout_initiated and not st.session_state.feedback_submitted:
        # Clear the page and show only feedback form
        st.title("📝 Session Feedback")
        st.write("Before you go, please provide some feedback about your session.")
        
        feedback_col1, feedback_col2, feedback_col3 = st.columns(3)
        
        with feedback_col1:
            topic = st.text_input("What topics did you discuss today?", key="logout_topic")
        
        with feedback_col2:
            rating = st.slider("How helpful was this session?", 1, 5, 3, key="logout_rating")
        
        with feedback_col3:
            difficulty = st.slider("How difficult were the topics?", 1, 5, 3, key="logout_difficulty")
        
        additional_feedback = st.text_area("Any additional comments or suggestions?", key="logout_comments")
        
        col1, col2 = st.columns([1, 5])
        with col1:
            if st.button("Submit Feedback"):
                if topic:
                    # Save final usage data
                    save_registration(st.session_state.user_data, st.session_state.start_time)
                    
                    # Save feedback
                    save_feedback(rating, topic, difficulty)
                    track_topic(topic, difficulty)
                    track_completion(True)
                    
                    # If there are additional comments, save them
                    if additional_feedback:
                        track_feedback_trend(rating, additional_feedback)
                    else:
                        track_feedback_trend(rating)
                    
                    st.session_state.feedback_submitted = True
                    st.rerun()
                else:
                    st.error("Please enter the topics discussed.")
            
            if st.button("Skip Feedback"):
                # Save final usage data without feedback
                save_registration(st.session_state.user_data, st.session_state.start_time)
                st.session_state.feedback_submitted = True
                st.rerun()
        
        # Stop here to prevent showing main app content
        st.stop()
    
    # Complete logout after feedback
    if st.session_state.logout_initiated and st.session_state.feedback_submitted:
        # Reset all session state variables
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        
        # Reinitialize essential variables
        st.session_state.registered = False
        st.session_state.chat_started = False
        st.session_state.logout_initiated = False
        st.session_state.feedback_submitted = False
        st.rerun()
    
    # Check if current time is within tutoring hours
    if is_within_tutoring_hours():
        # Show debug information in an expander
        with st.expander("Debug Time Information"):
            st.json(st.session_state.debug_time)
            
        st.warning("""
        ⚠️ In-person tutoring is currently available! 
        
        Please visit the in-person tutoring session instead of using the bot. 
        The bot will be available after the tutoring session ends.
        """)
        st.stop()

    # Get the API key from environment variable or secrets
    openai_api_key = os.environ.get("OPENAI_API_KEY") or st.secrets.get("OPENAI_API_KEY")

    if not openai_api_key:
        st.error("""
        ⚠️ OpenAI API Key not configured!
        
        To use NuAnswers, you need to:
        1. Get an API key from OpenAI (https://platform.openai.com/api-keys)
        2. Add it as an environment variable:
           - Key: OPENAI_API_KEY
           - Value: your-api-key-here
        3. Restart the application
        
        If you're deploying on Render:
        1. Go to your Render dashboard
        2. Select your service
        3. Go to the "Environment" tab
        4. Add a new environment variable:
           - Key: OPENAI_API_KEY
           - Value: your-api-key-here
        5. Redeploy your service
        """)
        st.stop()

    # Create an OpenAI client
    client = OpenAI(api_key=openai_api_key)

    # Move these two functions above the file upload section
    def encode_image_to_base64(file):
        """Convert uploaded image file to base64 string"""
        return base64.b64encode(file.getvalue()).decode('utf-8')

    def analyze_image(file):
        """Analyze image content using OpenAI's GPT-4 Vision model"""
        try:
            base64_image = encode_image_to_base64(file)
            response = client.chat.completions.create(
                model="gpt-4-vision-preview",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Please analyze this image in the context of accounting, finance, or business studies. Describe any relevant equations, problems, charts, or concepts shown."},
                            {
                                "type": "image_url",
                                "image_url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        ]
                    }
                ],
                max_tokens=300
            )
            return response.choices[0].message.content
        except Exception as e:
            st.error(f"Error analyzing image: {str(e)}")
            return None

    # File upload section
    st.subheader("📄 Upload Course Materials")
    uploaded_files = st.file_uploader(
        "Upload your course materials (PDF, DOCX, TXT, PPTX, CSV, XLS, XLSX, PNG, JPG, JPEG)",
        type=['pdf', 'docx', 'txt', 'pptx', 'csv', 'xls', 'xlsx', 'png', 'jpg', 'jpeg'],
        accept_multiple_files=True
    )
    
    if uploaded_files:
        for file in uploaded_files:
            if file not in [doc['file'] for doc in st.session_state.uploaded_documents]:
                file_extension = Path(file.name).suffix.lower()
                
                # Handle image files differently
                if file_extension in ['.png', '.jpg', '.jpeg']:
                    # Analyze image content
                    image_analysis = analyze_image(file)
                    
                    st.session_state.uploaded_documents.append({
                        'file': file,
                        'name': file.name,
                        'content': f"[Image Analysis: {image_analysis}]" if image_analysis else f"[Image File: {file.name}]",
                        'is_image': True,
                        'image_analysis': image_analysis
                    })
                    st.success(f"Successfully uploaded and analyzed image {file.name}")
                else:
                    text = extract_text_from_file(file)
                    if text:
                        st.session_state.uploaded_documents.append({
                            'file': file,
                            'name': file.name,
                            'content': text,
                            'is_image': False
                        })
                        st.success(f"Successfully processed {file.name}")
    
    # Search and document management section
    if st.session_state.uploaded_documents:
        st.subheader("📚 Your Uploaded Materials")
        
        # Search bar
        search_col, reorder_col = st.columns([3, 1])
        with search_col:
            st.session_state.search_query = st.text_input("🔍 Search in documents", 
                                                        value=st.session_state.search_query,
                                                        placeholder="Search by filename or content")
        
        # Reorder button
        with reorder_col:
            if st.button("🔄 Reorder Documents"):
                st.session_state.show_reorder = not getattr(st.session_state, 'show_reorder', False)
        
        # Reorder interface
        if getattr(st.session_state, 'show_reorder', False):
            st.info("Drag and drop documents to reorder them")
            for i, doc in enumerate(st.session_state.uploaded_documents):
                cols = st.columns([1, 4, 1])
                with cols[0]:
                    st.write(f"{i+1}.")
                with cols[1]:
                    st.write(doc['name'])
                with cols[2]:
                    if st.button("↑", key=f"up_{i}") and i > 0:
                        st.session_state.uploaded_documents[i], st.session_state.uploaded_documents[i-1] = \
                            st.session_state.uploaded_documents[i-1], st.session_state.uploaded_documents[i]
                        st.rerun()
                    if st.button("↓", key=f"down_{i}") and i < len(st.session_state.uploaded_documents) - 1:
                        st.session_state.uploaded_documents[i], st.session_state.uploaded_documents[i+1] = \
                            st.session_state.uploaded_documents[i+1], st.session_state.uploaded_documents[i]
                        st.rerun()
        
        # Display filtered documents
        filtered_docs = search_in_documents(st.session_state.search_query, st.session_state.uploaded_documents)
        
        if not filtered_docs:
            st.info("No documents match your search query.")
        else:
            for i, doc in enumerate(filtered_docs):
                cols = st.columns([4, 1])
                with cols[0].expander(doc['name']):
                    if doc.get('is_image', False):
                        st.image(doc['file'], caption=doc['name'])
                        if doc.get('image_analysis'):
                            st.markdown("**Image Analysis:**")
                            st.markdown(doc['image_analysis'])
                    else:
                        # Highlight search terms in content
                        content = doc['content']
                        if st.session_state.search_query:
                            query = st.session_state.search_query.lower()
                            start = content.lower().find(query)
                            if start != -1:
                                end = start + len(query)
                                highlighted = (
                                    content[:start] +
                                    f"**{content[start:end]}**" +
                                    content[end:]
                                )
                                st.markdown(highlighted)
                            else:
                                st.text(content[:500] + "..." if len(content) > 500 else content)
                        else:
                            st.text(content[:500] + "..." if len(content) > 500 else content)
                
                # Delete button with confirmation
                if cols[1].button("🗑️", key=f"delete_{i}"):
                    st.session_state.doc_to_delete = doc
                
                # Confirmation dialog
                if st.session_state.doc_to_delete == doc:
                    st.warning(f"Are you sure you want to delete {doc['name']}?")
                    confirm_cols = st.columns(2)
                    if confirm_cols[0].button("Yes, delete it", key=f"confirm_delete_{i}"):
                        st.session_state.uploaded_documents.remove(doc)
                        st.session_state.doc_to_delete = None
                        st.rerun()
                    if confirm_cols[1].button("Cancel", key=f"cancel_delete_{i}"):
                        st.session_state.doc_to_delete = None
                        st.rerun()

    # Display the existing chat messages
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Create a chat input field
    if prompt := st.chat_input("What would you like to work on today?"):
        response_start_time = datetime.now(ZoneInfo("America/New_York"))
        
        # Store and display the current prompt
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # Place this before the chat input logic (before the first use of track_response_time)
        def track_response_time(start_time, end_time):
            """Track system response time"""
            response_time = (end_time - start_time).total_seconds()
            entry = {
                "timestamp": end_time.strftime("%Y-%m-%d %H:%M:%S"),
                "response_time": response_time,
                "user_id": st.session_state.user_data.get("full_name")
            }
            st.session_state.response_times.append(entry)
            save_to_csv(entry, RESPONSE_TIMES_PATH)

        # Prepare context from uploaded documents
        context = ""
        if st.session_state.uploaded_documents:
            context = "\n\n".join([f"Document: {doc['name']}\nContent: {doc['content']}" 
                                 for doc in st.session_state.uploaded_documents])
            context = f"Here is the context from uploaded documents:\n\n{context}\n\n"

        # Generate a response using the OpenAI API
        stream = client.chat.completions.create(
            model="gpt-4.1",
            messages=[
                {"role": "system", "content": f"""You are an Accounting & Finance Tutor. Your role is to guide students through their homework and exam preparation through a conversational, step-by-step approach.

IMPORTANT RULES:
1. NEVER give direct answers or solutions
2. Ask ONE question at a time and wait for the student's response
3. After each student response, ask a follow-up question to guide their thinking
4. If the student's answer is incorrect, ask a guiding question to help them think differently
5. If the student asks for the answer, respond with a question that helps them think about the problem differently
6. Use simple, clear questions that build on each other
7. Focus on one concept or step at a time
8. Validate their understanding before moving to the next step
9. Use encouraging phrases like "Good thinking!" or "You're on the right track!"
10. If the student seems stuck, ask a simpler question that breaks down the problem
11. Use the context from uploaded documents to provide more relevant guidance

Example of good tutoring:
Student: "How do I solve this problem?"
Tutor: "Let's start with the first step. What information do we have in the problem?"
Student: [responds]
Tutor: "Good! Now, what do you think we should do with this information?"
[continue with one question at a time]

Example of bad tutoring:
"Here's how to solve it: First, do this, then do that, then calculate this..."
[giving multiple steps at once]"""},
                {"role": "system", "content": context},
                *[{"role": m["role"], "content": m["content"]} for m in st.session_state.messages]
            ],
            stream=True,
        )

        # Stream the response
        with st.chat_message("assistant"):
            response = st.write_stream(stream)
        st.session_state.messages.append({"role": "assistant", "content": response})

        response_end_time = datetime.now(ZoneInfo("America/New_York"))
        track_response_time(response_start_time, response_end_time)
        
        # Track content access if documents are referenced
        if st.session_state.uploaded_documents:
            try:
                for doc in st.session_state.uploaded_documents:
                    track_content_access(doc['name'], 'document' if not doc.get('is_image') else 'image')
            except NameError:
                pass  # track_content_access not in scope (e.g. some Streamlit run contexts)
        
        # Track resolution time if topic is completed
        if "current_topic_start" in st.session_state and "current_topic" in st.session_state:
            track_resolution_time(
                st.session_state.current_topic_start,
                datetime.now(ZoneInfo("America/New_York")),
                st.session_state.current_topic
            )

    # New chat: save current session and show course form again
    if st.sidebar.button("New chat"):
        save_registration(st.session_state.user_data, st.session_state.start_time)
        st.session_state.chat_started = False
        st.session_state.messages = [
            {"role": "assistant", "content": "Hello! I'm NuAnswers. I'm here to help you understand concepts and work through problems. What would you like to work on today?"}
        ]
        st.rerun()
    # Add a logout button at the top of the main content
    if st.sidebar.button("Logout"):
        st.session_state.logout_initiated = True
        st.rerun()

def show_admin_panel():
    """Display admin panel with registration statistics"""
    st.header("Admin Panel")
    
    try:
        # Load data from CSV
        csv_path = "registration_data.csv"
        if os.path.exists(csv_path):
            df = pd.read_csv(csv_path)
            
            # Convert timestamp to datetime
            df['timestamp'] = pd.to_datetime(df['timestamp'])
            
            # Calculate statistics
            total_registrations = len(df)
            total_usage_minutes = df['usage_time_minutes'].sum()
            avg_usage_minutes = df['usage_time_minutes'].mean()
            
            # Display statistics
            st.metric("Total Registrations", total_registrations)
            st.metric("Total Usage Time (minutes)", f"{total_usage_minutes:.2f}")
            st.metric("Average Usage Time (minutes)", f"{avg_usage_minutes:.2f}")
            
            # Show daily statistics
            st.subheader("Daily Statistics")
            daily_stats = df.groupby(df['timestamp'].dt.date).agg({
                'usage_time_minutes': ['sum', 'mean']
            }).reset_index()
            daily_stats.columns = ['Date', 'Total Minutes', 'Avg Minutes']
            st.dataframe(daily_stats.sort_values('Date', ascending=False))
            
            # Show raw data
            st.subheader("Raw Registration Data")
            st.dataframe(df.sort_values('timestamp', ascending=False))
        else:
            st.info("No registration data available yet.")
    except Exception as e:
        st.error(f"Error loading registration data: {str(e)}")

def save_feedback(rating, topic, difficulty):
    """Save feedback data"""
    et_tz = ZoneInfo("America/New_York")
    feedback_entry = {
        "timestamp": datetime.now(et_tz).strftime("%Y-%m-%d %H:%M:%S"),
        "course_id": st.session_state.user_data.get("course_id"),
        "rating": rating,
        "topic": topic,
        "difficulty": difficulty
    }
    st.session_state.feedback_data.append(feedback_entry)
    save_to_csv(feedback_entry, FEEDBACK_DATA_PATH)

def track_topic(topic, difficulty=None):
    """Track topic data"""
    et_tz = ZoneInfo("America/New_York")
    topic_entry = {
        "timestamp": datetime.now(et_tz).strftime("%Y-%m-%d %H:%M:%S"),
        "course_id": st.session_state.user_data.get("course_id"),
        "topic": topic,
        "difficulty": difficulty
    }
    st.session_state.topic_data.append(topic_entry)
    save_to_csv(topic_entry, TOPIC_DATA_PATH)

def track_completion(completed):
    """Track course completion"""
    et_tz = ZoneInfo("America/New_York")
    completion_entry = {
        "timestamp": datetime.now(et_tz).strftime("%Y-%m-%d %H:%M:%S"),
        "course_id": st.session_state.user_data.get("course_id"),
        "completed": completed
    }
    st.session_state.completion_data.append(completion_entry)
    save_to_csv(completion_entry, COMPLETION_DATA_PATH)

def track_system_status(status, start_time, end_time=None):
    """Track system uptime and status"""
    if end_time is None:
        end_time = datetime.now(ZoneInfo("America/New_York"))
    duration = (end_time - start_time).total_seconds()
    entry = {
        "start_time": start_time.strftime("%Y-%m-%d %H:%M:%S"),
        "end_time": end_time.strftime("%Y-%m-%d %H:%M:%S"),
        "status": status,
        "duration": duration
    }
    st.session_state.system_status.append(entry)
    save_to_csv(entry, SYSTEM_STATUS_PATH)

def track_feedback_trend(satisfaction_score, suggestions=None):
    """Track feedback trends over time"""
    entry = {
        "date": datetime.now(ZoneInfo("America/New_York")).strftime("%Y-%m-%d"),
        "satisfaction_score": satisfaction_score,
        "suggestions": suggestions,
        "user_id": st.session_state.user_data.get("full_name")
    }
    st.session_state.feedback_trends.append(entry)
    save_to_csv(entry, FEEDBACK_TRENDS_PATH)

def track_yearly_data():
    """Track yearly performance metrics"""
    current_year = datetime.now(ZoneInfo("America/New_York")).year
    entry = {
        "year": current_year,
        "registrations": len(st.session_state.registration_data),
        "unique_users": st.session_state.registration_data['full_name'].nunique(),
        "total_usage": st.session_state.registration_data['usage_time_minutes'].sum()
    }
    st.session_state.yearly_data.append(entry)
    save_to_csv(entry, YEARLY_DATA_PATH)

def get_current_semester():
    """
    Determine the current semester based on the academic calendar.
    Academic Calendar 2024-2025:
    - Fall 2024: 8/26/24 - 12/18/24
    - Winter 2025: 1/2/25 - 1/20/25
    - Spring 2025: 1/21/25 - 5/13/25
    - Summer 2025: 5/19/25 - 8/9/25
    """
    current_date = datetime.now(ZoneInfo("America/New_York"))
    
    # Define semester date ranges
    semester_dates = {
        ("2024-08-26", "2024-12-18"): ("Fall", "2024"),
        ("2025-01-02", "2025-01-20"): ("Winter", "2025"),
        ("2025-01-21", "2025-05-13"): ("Spring", "2025"),
        ("2025-05-19", "2025-08-09"): ("Summer", "2025")
    }
    
    current_date_str = current_date.strftime("%Y-%m-%d")
    
    for (start_date, end_date), (semester, year) in semester_dates.items():
        if start_date <= current_date_str <= end_date:
            return semester, year
    
    # If date falls between semesters, assign to the upcoming semester
    if current_date_str < "2024-08-26":
        return "Pre-Fall", "2024"
    elif "2024-12-18" < current_date_str < "2025-01-02":
        return "Winter-Break", "2024-2025"
    elif "2025-01-20" < current_date_str < "2025-01-21":
        return "Winter-Spring-Break", "2025"
    elif "2025-05-13" < current_date_str < "2025-05-19":
        return "Spring-Summer-Break", "2025"
    elif current_date_str > "2025-08-09":
        return "Post-Summer", "2025"
    
    return "Unknown", str(current_date.year)

def track_semester_data(semester=None, year=None):
    """Track semester-specific metrics with accurate semester classification"""
    if semester is None or year is None:
        semester, year = get_current_semester()
    
    # Get data for the current semester
    semester_start = {
        "Fall": "2024-08-26",
        "Winter": "2025-01-02",
        "Spring": "2025-01-21",
        "Summer": "2025-05-19"
    }.get(semester)
    
    semester_end = {
        "Fall": "2024-12-18",
        "Winter": "2025-01-20",
        "Spring": "2025-05-13",
        "Summer": "2025-08-09"
    }.get(semester)
    
    if semester_start and semester_end:
        # Filter registration data for the current semester
        mask = (st.session_state.registration_data['timestamp'] >= semester_start) & \
               (st.session_state.registration_data['timestamp'] <= semester_end)
        semester_data = st.session_state.registration_data[mask]
    else:
        semester_data = st.session_state.registration_data
    
    entry = {
        "year": year,
        "semester": semester,
        "start_date": semester_start,
        "end_date": semester_end,
        "registrations": len(semester_data),
        "unique_users": semester_data['full_name'].nunique() if not semester_data.empty else 0,
        "total_usage": semester_data['usage_time_minutes'].sum() if not semester_data.empty else 0,
        "avg_session_length": semester_data['usage_time_minutes'].mean() if not semester_data.empty else 0,
        "timestamp": datetime.now(ZoneInfo("America/New_York")).strftime("%Y-%m-%d %H:%M:%S")
    }
    
    st.session_state.semester_data.append(entry)
    save_to_csv(entry, SEMESTER_DATA_PATH)

def calculate_department_metrics(department_data):
    """Calculate comprehensive metrics for department performance"""
    metrics = {
        "total_students": department_data['full_name'].nunique(),
        "total_sessions": len(department_data),
        "total_usage_hours": department_data['usage_time_minutes'].sum() / 60,
        "avg_session_length": department_data['usage_time_minutes'].mean(),
        "avg_sessions_per_student": len(department_data) / department_data['full_name'].nunique() if department_data['full_name'].nunique() > 0 else 0,
    }
    
    # Calculate engagement metrics
    if not department_data.empty:
        # Update to use tuple for groupby
        repeat_users = department_data.groupby(('full_name',)).filter(lambda x: len(x) > 1)['full_name'].unique()
        metrics.update({
            "repeat_users": len(repeat_users),
            "repeat_user_rate": len(repeat_users) / department_data['full_name'].nunique() if department_data['full_name'].nunique() > 0 else 0,
        })
    
    return metrics

def track_department_data(department):
    """Enhanced department performance tracking"""
    # Get current semester info
    current_semester, current_year = get_current_semester()
    
    # Filter data for the department
    dept_data = st.session_state.registration_data[
        st.session_state.registration_data['major'] == department
    ]
    
    # Get semester dates
    semester_start = {
        "Fall": "2024-08-26",
        "Winter": "2025-01-02",
        "Spring": "2025-01-21",
        "Summer": "2025-05-19"
    }.get(current_semester)
    
    semester_end = {
        "Fall": "2024-12-18",
        "Winter": "2025-01-20",
        "Spring": "2025-05-13",
        "Summer": "2025-08-09"
    }.get(current_semester)
    
    # Filter for current semester if dates are available
    if semester_start and semester_end:
        mask = (dept_data['timestamp'] >= semester_start) & \
               (dept_data['timestamp'] <= semester_end)
        current_semester_data = dept_data[mask]
    else:
        current_semester_data = dept_data
    
    # Calculate usage patterns
    if not current_semester_data.empty:
        # Update to use tuples for groupby operations
        usage_patterns = {
            "peak_usage_hour": current_semester_data.groupby(('timestamp.hour',))['full_name'].count().idxmax()[0],
            "peak_usage_day": current_semester_data.groupby(('timestamp.day_name',))['full_name'].count().idxmax()[0],
            "avg_daily_users": current_semester_data.groupby(('timestamp.date',))['full_name'].nunique().mean()
        }
    else:
        usage_patterns = {
            "peak_usage_hour": None,
            "peak_usage_day": None,
            "avg_daily_users": 0
        }
    
    # Calculate current semester metrics
    current_metrics = calculate_department_metrics(current_semester_data)
    
    # Calculate satisfaction metrics if feedback data exists
    dept_feedback = [f for f in st.session_state.feedback_data 
                    if f.get('full_name') in dept_data['full_name'].unique()]
    
    satisfaction_metrics = {
        "avg_satisfaction": sum(f['rating'] for f in dept_feedback) / len(dept_feedback) if dept_feedback else 0,
        "total_feedback": len(dept_feedback),
        "satisfaction_rate_5": len([f for f in dept_feedback if f['rating'] == 5]) / len(dept_feedback) if dept_feedback else 0,
        "satisfaction_rate_4_plus": len([f for f in dept_feedback if f['rating'] >= 4]) / len(dept_feedback) if dept_feedback else 0
    }
    
    # Calculate topic performance
    dept_topics = [t for t in st.session_state.topic_data 
                  if t.get('full_name') in dept_data['full_name'].unique()]
    
    topic_metrics = {}
    if dept_topics:
        topic_counts = {}
        topic_difficulties = {}
        for t in dept_topics:
            topic = t.get('topic', '')
            difficulty = t.get('difficulty', 0)
            topic_counts[topic] = topic_counts.get(topic, 0) + 1
            topic_difficulties[topic] = topic_difficulties.get(topic, []) + [difficulty]
        
        topic_metrics = {
            "most_common_topics": sorted(topic_counts.items(), key=lambda x: x[1], reverse=True)[:5],
            "avg_topic_difficulty": {topic: sum(difficulties)/len(difficulties) 
                                   for topic, difficulties in topic_difficulties.items()}
        }
    
    # Calculate completion metrics
    dept_completions = [c for c in st.session_state.completion_data 
                       if c.get('full_name') in dept_data['full_name'].unique()]
    
    completion_metrics = {
        "completion_rate": len([c for c in dept_completions if c.get('completed', False)]) / len(dept_completions) if dept_completions else 0,
        "total_completions": len([c for c in dept_completions if c.get('completed', False)])
    }
    
    # Calculate usage patterns
    if not current_semester_data.empty:
        usage_patterns = {
            "peak_usage_hour": current_semester_data.groupby(('timestamp.hour',))['full_name'].count().idxmax()[0],
            "peak_usage_day": current_semester_data.groupby(('timestamp.day_name',))['full_name'].count().idxmax()[0],
            "avg_daily_users": current_semester_data.groupby(('timestamp.date',))['full_name'].nunique().mean()
        }
    else:
        usage_patterns = {
            "peak_usage_hour": None,
            "peak_usage_day": None,
            "avg_daily_users": 0
        }
    
    # Combine all metrics
    entry = {
        "timestamp": datetime.now(ZoneInfo("America/New_York")).strftime("%Y-%m-%d %H:%M:%S"),
        "department": department,
        "semester": current_semester,
        "year": current_year,
        **current_metrics,
        **satisfaction_metrics,
        **completion_metrics,
        **usage_patterns,
        "topic_analysis": topic_metrics
    }
    
    # Calculate relative performance (compared to other departments)
    all_dept_data = {}
    for dept in st.session_state.registration_data['major'].unique():
        other_dept_data = st.session_state.registration_data[
            st.session_state.registration_data['major'] == dept
        ]
        all_dept_data[dept] = calculate_department_metrics(other_dept_data)
    
    # Add relative performance metrics
    if len(all_dept_data) > 1:
        dept_metrics = all_dept_data[department]
        other_depts_avg = {
            metric: sum(d[metric] for d in all_dept_data.values() if d != dept_metrics) / (len(all_dept_data) - 1)
            for metric in dept_metrics.keys()
        }
        
        relative_performance = {
            f"relative_{metric}": (dept_metrics[metric] / other_depts_avg[metric] if other_depts_avg[metric] > 0 else 1)
            for metric in dept_metrics.keys()
        }
        
        entry.update(relative_performance)
    
    st.session_state.department_data.append(entry)
    save_to_csv(entry, DEPARTMENT_DATA_PATH)
    
    return entry

def track_historical_usage():
    """Track historical usage patterns"""
    entry = {
        "date": datetime.now(ZoneInfo("America/New_York")).strftime("%Y-%m-%d"),
        "usage": st.session_state.registration_data['usage_time_minutes'].sum() / 60,  # Convert to hours
        "unique_users": st.session_state.registration_data['full_name'].nunique()
    }
    st.session_state.historical_usage.append(entry)
    save_to_csv(entry, HISTORICAL_USAGE_PATH)

def track_hourly_usage():
    """Track hourly usage patterns"""
    current_hour = datetime.now(ZoneInfo("America/New_York")).hour
    entry = {
        "timestamp": datetime.now(ZoneInfo("America/New_York")).strftime("%Y-%m-%d %H:%M:%S"),
        "hour": current_hour,
        "usage": len(st.session_state.registration_data[
            st.session_state.registration_data['timestamp'].dt.hour == current_hour
        ])
    }
    st.session_state.hourly_usage.append(entry)
    save_to_csv(entry, HOURLY_USAGE_PATH)

def calculate_success_indicators(student_data):
    """Calculate success indicators for a student based on their usage patterns"""
    if student_data.empty:
        return {}
    
    # Usage patterns
    total_sessions = len(student_data)
    total_hours = student_data['usage_time_minutes'].sum() / 60
    avg_session_length = student_data['usage_time_minutes'].mean()
    
    # Engagement patterns
    days_active = student_data['timestamp'].dt.date.nunique()
    sessions_per_day = total_sessions / days_active if days_active > 0 else 0
    
    # Calculate consistency score (0-1)
    date_range = (student_data['timestamp'].max() - student_data['timestamp'].min()).days + 1
    consistency_score = days_active / date_range if date_range > 0 else 0
    
    return {
        "total_sessions": total_sessions,
        "total_hours": total_hours,
        "avg_session_length": avg_session_length,
        "days_active": days_active,
        "sessions_per_day": sessions_per_day,
        "consistency_score": consistency_score
    }

def calculate_topic_mastery(student_id):
    """Calculate topic mastery levels for a student"""
    student_topics = [t for t in st.session_state.topic_data if t.get('full_name') == student_id]
    
    topic_mastery = {}
    for topic_data in student_topics:
        topic = topic_data.get('topic', '')
        difficulty = topic_data.get('difficulty', 0)
        
        if topic not in topic_mastery:
            topic_mastery[topic] = {
                'attempts': 0,
                'total_difficulty': 0,
                'completions': 0
            }
        
        topic_mastery[topic]['attempts'] += 1
        topic_mastery[topic]['total_difficulty'] += difficulty
        
        # Check if topic was completed successfully
        completion_records = [c for c in st.session_state.completion_data 
                            if c.get('full_name') == student_id and 
                            c.get('topic') == topic and 
                            c.get('completed', False)]
        topic_mastery[topic]['completions'] += len(completion_records)
    
    # Calculate mastery scores
    mastery_scores = {}
    for topic, data in topic_mastery.items():
        attempts = data['attempts']
        avg_difficulty = data['total_difficulty'] / attempts if attempts > 0 else 0
        completion_rate = data['completions'] / attempts if attempts > 0 else 0
        
        # Mastery score (0-1) based on completion rate and difficulty
        mastery_scores[topic] = (completion_rate * 0.7 + (avg_difficulty / 5) * 0.3)
    
    return mastery_scores

def predict_student_success(student_id):
    """Predict student success based on various metrics"""
    # Get student data
    student_data = st.session_state.registration_data[
        st.session_state.registration_data['full_name'] == student_id
    ]
    
    # Get current semester
    current_semester, current_year = get_current_semester()
    
    # Calculate base success indicators
    success_indicators = calculate_success_indicators(student_data)
    
    # Calculate topic mastery
    topic_mastery = calculate_topic_mastery(student_id)
    
    # Get feedback and satisfaction data
    student_feedback = [f for f in st.session_state.feedback_data if f.get('full_name') == student_id]
    avg_satisfaction = sum(f['rating'] for f in student_feedback) / len(student_feedback) if student_feedback else 0
    
    # Calculate engagement trend
    if not student_data.empty:
        student_data = student_data.sort_values('timestamp')
        # Update to use tuple for groupby
        weekly_sessions = student_data.groupby(pd.Grouper(key='timestamp', freq='W'))['full_name'].count()
        engagement_trend = weekly_sessions.pct_change().mean() if len(weekly_sessions) > 1 else 0
    else:
        engagement_trend = 0
    
    # Calculate success probability
    success_factors = {
        'usage_score': min(success_indicators.get('total_hours', 0) / 10, 1),  # Cap at 10 hours
        'consistency_score': success_indicators.get('consistency_score', 0),
        'mastery_score': sum(topic_mastery.values()) / len(topic_mastery) if topic_mastery else 0,
        'satisfaction_score': avg_satisfaction / 5 if avg_satisfaction > 0 else 0,
        'engagement_trend': (engagement_trend + 1) / 2 if engagement_trend > -1 else 0  # Normalize to 0-1
    }
    
    # Weight the factors
    weights = {
        'usage_score': 0.25,
        'consistency_score': 0.25,
        'mastery_score': 0.2,
        'satisfaction_score': 0.15,
        'engagement_trend': 0.15
    }
    
    # Calculate overall success probability
    success_probability = sum(score * weights[factor] for factor, score in success_factors.items())
    
    # Generate recommendations
    recommendations = []
    if success_factors['usage_score'] < 0.6:
        recommendations.append("Increase total usage time")
    if success_factors['consistency_score'] < 0.6:
        recommendations.append("Maintain more regular study sessions")
    if success_factors['mastery_score'] < 0.6:
        recommendations.append("Focus on completing more topics")
    if success_factors['satisfaction_score'] < 0.6:
        recommendations.append("Engage more actively with the content")
    if success_factors['engagement_trend'] < 0.5:
        recommendations.append("Increase weekly engagement")
    
    # Identify areas of strength
    strengths = [factor for factor, score in success_factors.items() if score >= 0.8]
    
    # Calculate risk level
    risk_level = "Low" if success_probability >= 0.7 else "Medium" if success_probability >= 0.4 else "High"
    
    # Prepare prediction results
    prediction = {
        "timestamp": datetime.now(ZoneInfo("America/New_York")).strftime("%Y-%m-%d %H:%M:%S"),
        "full_name": student_id,
        "semester": current_semester,
        "year": current_year,
        "success_probability": success_probability,
        "risk_level": risk_level,
        "success_factors": success_factors,
        "topic_mastery": topic_mastery,
        "usage_metrics": success_indicators,
        "recommendations": recommendations,
        "strengths": strengths,
        "avg_satisfaction": avg_satisfaction,
        "engagement_trend": engagement_trend
    }
    
    # Save prediction to session state and CSV
    if "success_predictions" not in st.session_state:
        st.session_state.success_predictions = []
    
    st.session_state.success_predictions.append(prediction)
    save_to_csv(prediction, STUDENT_PERFORMANCE_PATH)
    
    return prediction

def track_student_performance(student_id, usage_hours, success_rate):
    """Enhanced student performance tracking with success prediction"""
    # Calculate usage category
    usage_category = pd.cut([usage_hours], bins=[0, 5, 10, 15, float('inf')], 
                          labels=['Low', 'Medium', 'High', 'Very High'])[0]
    
    # Get success prediction
    prediction = predict_student_success(student_id)
    
    # Combine metrics
    entry = {
        "timestamp": datetime.now(ZoneInfo("America/New_York")).strftime("%Y-%m-%d %H:%M:%S"),
        "full_name": student_id,
        "usage_hours": usage_hours,
        "success_rate": success_rate,
        "usage_category": usage_category,
        "success_probability": prediction['success_probability'],
        "risk_level": prediction['risk_level'],
        "recommendations": prediction['recommendations']
    }
    
    st.session_state.student_performance.append(entry)
    save_to_csv(entry, STUDENT_PERFORMANCE_PATH)
    
    return entry
